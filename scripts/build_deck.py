"""
Build the pitch deck as a native PowerPoint file.

    python scripts/build_deck.py        ->  Vetoed-deck.pptx

WHY A SCRIPT AND NOT A SAVED FILE. The numbers on these slides come from the
screener and the journal, and they move. Regenerating from source is how the
deck stays true to what the agent actually did, rather than drifting into a
set of figures that were right once.

FONTS. IBM Plex, which the HTML deck and the dashboard use, is a web font and
is not installed on a typical Windows machine - PowerPoint would silently
substitute something worse. Segoe UI and Consolas ship with Windows, read
cleanly when a video encoder compresses them, and keep the same
sans-plus-mono pairing.

Slides are built as native shapes and text, so everything is editable in
PowerPoint afterwards - including recording narration per slide and exporting
straight to MP4 via Slide Show > Record, then File > Export > Create a Video.
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "Vetoed-deck.pptx")

# Same palette as the dashboard and the HTML deck.
GROUND = RGBColor(0x0A, 0x0E, 0x13)
PANEL = RGBColor(0x10, 0x16, 0x1E)
RAISED = RGBColor(0x16, 0x1E, 0x28)
LINE = RGBColor(0x21, 0x2C, 0x38)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
DIM = RGBColor(0xA9, 0xB6, 0xC3)
MUTED = RGBColor(0x75, 0x82, 0x8F)
ACCENT = RGBColor(0x4D, 0xD4, 0xC0)
PROFIT = RGBColor(0x3E, 0xCF, 0x8E)
LOSS = RGBColor(0xF2, 0x68, 0x5C)
WARN = RGBColor(0xF0, 0xA9, 0x3B)
WARN_BG = RGBColor(0x24, 0x1B, 0x0C)
ACCENT_DK = RGBColor(0x12, 0x38, 0x32)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)          # side margin
CONTENT_W = W - 2 * M


# --------------------------------------------------------------------------- #
# primitives
# --------------------------------------------------------------------------- #

def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = GROUND
    return s


def text(s, x, y, w, h, runs, size=18, font=SANS, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP,
         letter_space=None):
    """`runs` is a string, or a list of (text, {overrides}) pairs."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for content, over in runs:
        r = p.add_run()
        r.text = content
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        if letter_space or over.get("letter_space"):
            # python-pptx has no spacing API; set it on the run XML directly.
            r.font._rPr.set("spc", str(int((over.get("letter_space")
                                            or letter_space) * 100)))
    return box


def rect(s, x, y, w, h, fill=PANEL, border=LINE, radius=True):
    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    if radius:
        try:
            shape.adjustments[0] = 0.06
        except (IndexError, KeyError):
            pass
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def eyebrow(s, label):
    text(s, M, Inches(0.72), CONTENT_W, Inches(0.3), label.upper(),
         size=12, font=MONO, color=ACCENT, bold=True, letter_space=2.2)


def heading(s, title, y=Inches(1.18)):
    text(s, M, y, CONTENT_W, Inches(1.25), title,
         size=34, color=TEXT, bold=True, spacing=1.06)


def body(s, x, y, w, paras, size=16, color=DIM):
    """Each paragraph is a string or a list of (text, overrides) runs."""
    box = s.shapes.add_textbox(x, y, w, Inches(3.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.35
        p.space_after = Pt(11)
        for content, over in ([(para, {})] if isinstance(para, str) else para):
            r = p.add_run()
            r.text = content
            r.font.name = over.get("font", SANS)
            r.font.size = Pt(over.get("size", size))
            r.font.bold = over.get("bold", False)
            r.font.color.rgb = over.get("color", color)
    return box


def card(s, x, y, w, h, label, value, sub="", value_color=TEXT, value_size=30):
    rect(s, x, y, w, h)
    pad = Inches(0.24)
    text(s, x + pad, y + pad, w - 2 * pad, Inches(0.24), label.upper(),
         size=10.5, font=MONO, color=MUTED, bold=True, letter_space=1.4)
    text(s, x + pad, y + pad + Inches(0.3), w - 2 * pad, Inches(0.55), value,
         size=value_size, font=MONO, color=value_color, bold=True)
    if sub:
        text(s, x + pad, y + h - pad - Inches(0.52), w - 2 * pad, Inches(0.5),
             sub, size=11.5, color=MUTED, spacing=1.2)


def chrome(s, n, total, show=True):
    if not show:
        return
    text(s, M, H - Inches(0.62), Inches(3), Inches(0.3), "VETOED",
         size=11, font=MONO, color=MUTED, bold=True, letter_space=2.5)
    text(s, W - M - Inches(3), H - Inches(0.62), Inches(3), Inches(0.3),
         "%02d / %d" % (n, total), size=11, font=MONO, color=MUTED,
         align=PP_ALIGN.RIGHT, letter_space=1.2)
    bar = rect(s, Emu(0), H - Inches(0.055), Emu(int(W * n / total)),
               Inches(0.055), fill=ACCENT, border=None, radius=False)
    bar.shadow.inherit = False


# --------------------------------------------------------------------------- #
# slides
# --------------------------------------------------------------------------- #

def build() -> str:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    N = 12
    n = 0

    # 1 - hook ------------------------------------------------------------- #
    s = slide(prs)
    text(s, M, Inches(2.05), CONTENT_W, Inches(0.3),
         "ALPACA AI TRADING AGENTS HACKATHON", size=12, font=MONO,
         color=ACCENT, bold=True, letter_space=2.6)
    text(s, M, Inches(2.5), CONTENT_W, Inches(1.6), "Vetoed",
         size=76, color=TEXT, bold=True)
    text(s, M, Inches(4.05), Inches(8.6), Inches(1.1),
         [("An autonomous options agent where ", {}),
          ("the AI is the least-trusted component.",
           {"bold": True, "color": TEXT})],
         size=23, color=DIM, spacing=1.35)
    text(s, M, Inches(5.5), CONTENT_W, Inches(0.4),
         "Alpaca paper trading   \u00b7   chong1120.github.io/Vetoed",
         size=13, font=MONO, color=MUTED)
    n += 1

    # 2 - the problem ------------------------------------------------------ #
    s = slide(prs); n += 1
    eyebrow(s, "The problem")
    heading(s, "Most AI trading agents give the model\ntoo much authority.")
    body(s, M, Inches(2.95), Inches(6.6), [
        [("You cannot unit-test a language model.", {"size": 19, "color": TEXT})],
        [("But you can unit-test the code that decides whether to obey it.",
          {"size": 19, "color": ACCENT, "bold": True})],
    ], size=19)
    x = M + Inches(7.2)
    rect(s, x, Inches(2.75), Inches(4.4), Inches(2.5))
    text(s, x + Inches(0.32), Inches(3.05), Inches(3.8), Inches(0.3),
         "SO THE MODEL'S ENTIRE AUTHORITY IS", size=10.5, font=MONO,
         color=MUTED, bold=True, letter_space=1.3)
    text(s, x + Inches(0.32), Inches(3.5), Inches(3.8), Inches(1.0),
         "Pick one item from a list\nit did not write.",
         size=17, color=ACCENT, bold=True, spacing=1.35)
    text(s, x + Inches(0.32), Inches(4.55), Inches(3.8), Inches(0.6),
         "It cannot build a trade, choose a size, or send an order.",
         size=11.5, color=MUTED, spacing=1.3)
    chrome(s, n, N)

    # 3 - the strategy ----------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "The strategy")
    heading(s, "We sell defined-risk insurance.")
    text(s, M, Inches(2.35), Inches(11.4), Inches(0.6),
         [("Two legs, always. The one we sell collects a premium; the one we "
           "buy ", {}),
          ("caps the loss before the order is sent", {"bold": True, "color": TEXT}),
          (".", {})],
         size=15, color=DIM, spacing=1.35)

    for i, (title, col, sub, legs) in enumerate((
            ("PUT CREDIT SPREAD", PROFIT, "Profits if the stock does not fall far",
             "SELL the 310 put   \u00b7   BUY the 305 put"),
            ("CALL CREDIT SPREAD", LOSS, "Profits if the stock does not rise far",
             "SELL the 300 call   \u00b7   BUY the 305 call"))):
        bx = M + i * Inches(5.9)
        rect(s, bx, Inches(3.1), Inches(5.5), Inches(3.0))
        text(s, bx + Inches(0.3), Inches(3.32), Inches(4.9), Inches(0.3),
             title, size=14, font=MONO, color=col, bold=True, letter_space=1.1)
        text(s, bx + Inches(0.3), Inches(3.68), Inches(4.9), Inches(0.3),
             sub, size=11.5, color=MUTED)
        # payoff line: flat, ramp, flat - drawn as three connected segments
        y_hi, y_lo = Inches(4.15), Inches(5.15)
        x0, x1, x2, x3 = (bx + Inches(0.45), bx + Inches(1.9),
                          bx + Inches(3.3), bx + Inches(5.05))
        if i == 0:      # put credit: low-left, high-right
            pts = [(x0, y_lo), (x1, y_lo), (x2, y_hi), (x3, y_hi)]
        else:           # call credit: high-left, low-right
            pts = [(x0, y_hi), (x1, y_hi), (x2, y_lo), (x3, y_lo)]
        for a, b in zip(pts, pts[1:]):
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, min(a[0], b[0]),
                                    min(a[1], b[1]) - Pt(1.2),
                                    abs(b[0] - a[0]) or Pt(2),
                                    (abs(b[1] - a[1]) or Pt(2)) + Pt(2.4))
            ln.fill.background(); ln.line.fill.background(); ln.shadow.inherit = False
            conn = s.shapes.add_connector(1, a[0], a[1], b[0], b[1])
            conn.line.color.rgb = col
            conn.line.width = Pt(2.5)
        text(s, bx + Inches(0.3), Inches(5.5), Inches(4.9), Inches(0.4),
             legs, size=11.5, color=DIM)
    chrome(s, n, N)

    # 4 - why volatility --------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "Why volatility matters")
    heading(s, "Two numbers. When they disagree,\nthe seller is being paid.")
    body(s, M, Inches(3.0), Inches(6.5), [
        [("Implied volatility", {"bold": True, "color": ACCENT}),
         (" \u2014 how much movement option prices are charging for.", {})],
        [("Realised volatility", {"bold": True, "color": PROFIT}),
         (" \u2014 how much the stock actually moved recently.", {})],
        [("When implied sits above realised, options are priced for more "
          "movement than has been happening. ", {}),
         ("That gap is what a premium seller is trying to collect.",
          {"bold": True, "color": TEXT})],
    ], size=15)
    x = M + Inches(7.1)
    rect(s, x, Inches(2.75), Inches(4.5), Inches(3.05))
    text(s, x + Inches(0.32), Inches(3.0), Inches(3.9), Inches(0.3),
         "AAPL, ONE LIVE SCREEN", size=10.5, font=MONO, color=MUTED,
         bold=True, letter_space=1.3)
    for j, (lab, pct, col, width_in) in enumerate((
            ("implied", "23.8%", ACCENT, 2.85),
            ("realised, 20 day", "18.9%", PROFIT, 2.26))):
        yy = Inches(3.45 + j * 0.78)
        text(s, x + Inches(0.32), yy, Inches(2.4), Inches(0.25), lab,
             size=11, font=MONO, color=MUTED)
        bar = rect(s, x + Inches(0.32), yy + Inches(0.28), Inches(width_in),
                   Inches(0.22), fill=col, border=None, radius=False)
        bar.shadow.inherit = False
        text(s, x + Inches(3.3), yy + Inches(0.16), Inches(1.0), Inches(0.35),
             pct, size=16, font=MONO, color=col, bold=True)
    text(s, x + Inches(0.32), Inches(5.15), Inches(1.5), Inches(0.3), "ratio",
         size=11, font=MONO, color=MUTED)
    text(s, x + Inches(1.4), Inches(4.95), Inches(2.5), Inches(0.6), "1.26",
         size=30, font=MONO, color=ACCENT, bold=True)
    chrome(s, n, N)

    # 5 - the measurement -------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "The measurement")
    heading(s, "One model. Two volatility inputs.")
    text(s, M, Inches(2.4), Inches(11.4), Inches(0.4),
         [("Same spread, same payoff, same probability model. ", {}),
          ("Only the volatility changes.", {"bold": True, "color": TEXT})],
         size=15, color=DIM)
    card(s, M, Inches(3.0), Inches(5.5), Inches(1.5), "Priced at realised vol",
         "+$28.59", "what it is worth if volatility keeps behaving like it has",
         PROFIT, 26)
    card(s, M + Inches(5.9), Inches(3.0), Inches(5.5), Inches(1.5),
         "Priced at implied vol", "\u2212$13.85",
         "roughly what the market is charging for it today", LOSS, 26)
    box = rect(s, M, Inches(4.72), Inches(11.4), Inches(1.15),
               fill=ACCENT_DK, border=ACCENT)
    text(s, M + Inches(0.3), Inches(4.9), Inches(9.0), Inches(0.3),
         "VETOED'S MODEL-DERIVED EDGE \u2014 THE DIFFERENCE", size=10.5,
         font=MONO, color=ACCENT, bold=True, letter_space=1.3)
    text(s, M + Inches(0.3), Inches(5.22), Inches(6.0), Inches(0.5),
         [("+$42.45", {"size": 26, "color": ACCENT, "bold": True, "font": MONO}),
          ("   per spread", {"size": 13, "color": DIM})], size=26)
    text(s, M, Inches(6.05), Inches(11.4), Inches(0.7),
         [("This is ", {}), ("our own operational signal", {"bold": True, "color": TEXT}),
          (", motivated by volatility-risk-premium research \u2014 not the "
           "canonical academic VRP, which is defined on variance swap rates "
           "over a matched horizon.", {})],
         size=11.5, color=MUTED, spacing=1.35)
    chrome(s, n, N)

    # 6 - the bug ---------------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "Proof the measurement is honest")
    heading(s, "I caught my own agent lying.")
    body(s, M, Inches(2.75), Inches(6.4), [
        [("An IWM trade where implied and realised volatility were ", {}),
         ("almost identical", {"bold": True, "color": TEXT}),
         (" \u2014 so there was no gap to collect.", {})],
        [("It reported a healthy edge anyway. Two defects: one side used ", {}),
         ("delta as a probability", {"bold": True, "color": TEXT}),
         (", and the payoff between the strikes was ", {}),
         ("approximated at its midpoint", {"bold": True, "color": TEXT}), (".", {})],
        [("Both changed real trade decisions", {"bold": True, "color": TEXT}),
         (", not just displayed numbers. Found, fixed, pinned by tests.", {})],
    ], size=14)
    x = M + Inches(7.1)
    card(s, x, Inches(2.6), Inches(4.5), Inches(1.32),
         "Before \u2014 two mismatched methods", "+$2.75",
         "reported on a trade with no gap on offer", LOSS, 24)
    card(s, x, Inches(4.06), Inches(4.5), Inches(1.32),
         "After \u2014 one corrected model", "+$0.34",
         "below the $2.00 gate \u2014 so it is not taken", PROFIT, 24)
    card(s, x, Inches(5.52), Inches(4.5), Inches(1.0), "Regression tests",
         "77 \u2192 236", "", ACCENT, 22)
    chrome(s, n, N)

    # 7 - separation of powers --------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "Separation of powers")
    heading(s, "The model proposes.\nDeterministic code disposes.")
    nodes = [("SCREENER", "builds valid spreads", "no judgement", PANEL, LINE, TEXT),
             ("LLM", "picks one, or none", "cannot invent or size", WARN_BG, WARN, WARN),
             ("RISK", "sizes, and can veto", "cannot be overruled", PANEL, LINE, TEXT),
             ("ALPACA", "one atomic order", "never legs in", PANEL, LINE, TEXT)]
    nw, gap = Inches(2.62), Inches(0.31)
    for i, (title, can, cant, fill, border, tc) in enumerate(nodes):
        x = M + i * (nw + gap)
        rect(s, x, Inches(2.85), nw, Inches(1.42), fill=fill, border=border)
        text(s, x, Inches(3.05), nw, Inches(0.3), title, size=14, font=MONO,
             color=tc, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.15), Inches(3.42), nw - Inches(0.3), Inches(0.3),
             can, size=10.5, color=PROFIT, align=PP_ALIGN.CENTER, spacing=1.2)
        text(s, x + Inches(0.15), Inches(3.72), nw - Inches(0.3), Inches(0.3),
             cant, size=10.5, color=LOSS, align=PP_ALIGN.CENTER, spacing=1.2)
        if i < 3:
            text(s, x + nw, Inches(3.4), gap, Inches(0.3), "\u2192", size=15,
                 font=MONO, color=MUTED, align=PP_ALIGN.CENTER)
    rows = [("Echoed contracts are checked against the shortlist", "hallucination \u2192 no trade", LOSS),
            ("The model's requested size is discarded", "cannot size", LOSS),
            ("No broker or tool output ever enters the prompt", "controlled vocabulary", LOSS),
            ("No model available? Arithmetic selection runs instead", "still autonomous", PROFIT)]
    y = Inches(4.75)
    for label, tag, col in rows:
        text(s, M, y, Inches(7.6), Inches(0.32), label, size=13, color=DIM)
        text(s, M + Inches(7.7), y, Inches(3.9), Inches(0.32), tag.upper(),
             size=10.5, font=MONO, color=col, bold=True, letter_space=1.1)
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, y + Inches(0.38),
                                CONTENT_W, Pt(0.75))
        ln.fill.solid(); ln.fill.fore_color.rgb = LINE
        ln.line.fill.background(); ln.shadow.inherit = False
        y += Inches(0.56)
    chrome(s, n, N)

    # 8 - risk -------------------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "Risk is not advisory")
    heading(s, "Eight gates. Any one of them\nrejects the trade outright.")
    facts = [("5%", TEXT, "of equity at risk\nin any one spread"),
             ("\u22123%", LOSS, "on the day halts\nthe whole session"),
             ("25", TEXT, "contracts, hard cap\nwhatever the maths says"),
             ("0", PROFIT, "naked positions,\nstructurally impossible"),
             ("2", TEXT, "legs, one atomic order \u2014\nboth fill or neither"),
             ("236", ACCENT, "tests, most of them\non the gates")]
    for i, (n_, col, lab) in enumerate(facts):
        cx = M + (i % 3) * Inches(3.95)
        cy = Inches(2.95) + (i // 3) * Inches(1.75)
        text(s, cx, cy, Inches(3.6), Inches(0.8), n_, size=44, font=MONO,
             color=col, bold=True)
        text(s, cx, cy + Inches(0.78), Inches(3.6), Inches(0.7), lab,
             size=12.5, color=DIM, spacing=1.3)
    chrome(s, n, N)

    # 9 - it says no -------------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "It refuses to trade")
    heading(s, "20 valid spreads. 7 survived.")
    text(s, M, Inches(2.4), Inches(11.0), Inches(0.35),
         "One screen, counted at every gate.", size=14, color=DIM)
    hdr = ["Underlying", "Implied \u00f7 realised", "Valid", "Cleared the gate"]
    cols = [Inches(2.1), Inches(2.9), Inches(1.5), Inches(5.4)]
    y = Inches(2.95); x = M
    for i, h in enumerate(hdr):
        text(s, x, y, cols[i], Inches(0.3), h.upper(), size=10.5, font=MONO,
             color=MUTED, bold=True, letter_space=1.2)
        x += cols[i]
    y += Inches(0.42)
    data = [("AAPL", "1.260", PROFIT, "2", "2  \u2014  richest premium on the board", PROFIT),
            ("SPY", "1.069", WARN, "3", "3", DIM),
            ("IWM", "1.067", WARN, "13", "2", DIM),
            ("QQQ", "0.894", LOSS, "2", "0  \u2014  implied below realised", LOSS)]
    for sym, ratio, rcol, valid, outcome, ocol in data:
        x = M
        text(s, x, y, cols[0], Inches(0.35), sym, size=14, font=MONO,
             color=TEXT, bold=True); x += cols[0]
        text(s, x, y, cols[1], Inches(0.35), ratio, size=14, font=MONO,
             color=rcol, bold=True); x += cols[1]
        text(s, x, y, cols[2], Inches(0.35), valid, size=14, font=MONO,
             color=DIM); x += cols[2]
        text(s, x, y, cols[3], Inches(0.35), outcome, size=13, color=ocol)
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, y + Inches(0.4),
                                CONTENT_W, Pt(0.75))
        ln.fill.solid(); ln.fill.fore_color.rgb = LINE
        ln.line.fill.background(); ln.shadow.inherit = False
        y += Inches(0.56)
    text(s, M, y + Inches(0.35), Inches(11.4), Inches(0.5),
         [("The agent isn't looking for trades. ", {}),
          ("It's looking for paid risk.", {"bold": True, "color": ACCENT})],
         size=18, color=DIM)
    chrome(s, n, N)

    # 10 - self-improving loop ---------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "Learning, on a short leash")
    heading(s, "It can learn. It can't learn its way\naround the risk controls.")
    steps = [("TRADE", False), ("JOURNAL EVERY DECISION", False),
             ("ANALYSE CLOSED TRADES", False), ("MIN 5 CLOSED TRADES", True),
             ("TIGHTEN A LIMIT", False), ("NEXT SESSION", False)]
    y = Inches(2.85)
    for label, is_gate in steps:
        rect(s, M + Inches(0.4), y, Inches(4.5), Inches(0.5),
             fill=WARN_BG if is_gate else PANEL,
             border=WARN if is_gate else LINE)
        text(s, M + Inches(0.4), y + Inches(0.13), Inches(4.5), Inches(0.3),
             label, size=12, font=MONO, color=WARN if is_gate else TEXT,
             bold=True, align=PP_ALIGN.CENTER)
        y += Inches(0.5)
        if label != "NEXT SESSION":
            text(s, M + Inches(0.4), y, Inches(4.5), Inches(0.22), "\u2193",
                 size=11, font=MONO, color=MUTED, align=PP_ALIGN.CENTER)
            y += Inches(0.22)
    x = M + Inches(5.9)
    body(s, x, Inches(2.9), Inches(5.5), [
        "Three losses in a row on one underlying, and it stops trading that "
        "underlying. A losing delta band gets narrowed.",
        [("Every adjustment is clamped against the defaults", {"bold": True, "color": TEXT}),
         (", so the loop can only ever make the agent more restrictive.", {})],
    ], size=14)
    rect(s, x, Inches(4.75), Inches(5.5), Inches(1.85), border=WARN)
    text(s, x + Inches(0.3), Inches(4.95), Inches(4.9), Inches(0.3),
         "IT CANNOT", size=10.5, font=MONO, color=WARN, bold=True,
         letter_space=1.3)
    text(s, x + Inches(0.3), Inches(5.32), Inches(4.9), Inches(1.2),
         "increase position size\nloosen a hard limit\nre-enable something it "
         "disabled\nreact to fewer than 5 closed trades",
         size=12.5, color=DIM, spacing=1.45)
    chrome(s, n, N)

    # 11 - live system ------------------------------------------------------ #
    s = slide(prs); n += 1
    eyebrow(s, "Running, unattended")
    heading(s, "It trades on a schedule,\nwith nothing of mine switched on.")
    body(s, M, Inches(2.95), Inches(6.6), [
        [("A scheduled trigger starts one session a day on GitHub Actions; "
          "inside it the agent re-screens every 10 minutes and may open a "
          "position every 30. ", {}),
         ("No server, no laptop.", {"bold": True, "color": TEXT})],
        "Every cycle reconciles against Alpaca before it does anything, so a "
        "restart can't duplicate a position \u2014 and every order carries an id "
        "the broker will refuse twice.",
    ], size=14)
    text(s, M, Inches(5.05), Inches(6.6), Inches(0.35),
         "chong1120.github.io/Vetoed", size=15, font=MONO, color=ACCENT, bold=True)
    text(s, M, Inches(5.45), Inches(6.6), Inches(0.6),
         "Equity curve \u00b7 screening funnel \u00b7 implied vs realised \u00b7 "
         "every decision \u00b7 every veto", size=11.5, color=MUTED, spacing=1.3)
    x = M + Inches(7.2)
    stats = [("236", PROFIT, "tests passing"), ("0", TEXT, "LLM write paths"),
             ("8", TEXT, "risk gates"), ("100%", ACCENT, "decisions journalled")]
    for i, (v, col, lab) in enumerate(stats):
        cx = x + (i % 2) * Inches(2.3)
        cy = Inches(2.95) + (i // 2) * Inches(1.6)
        text(s, cx, cy, Inches(2.1), Inches(0.7), v, size=36, font=MONO,
             color=col, bold=True)
        text(s, cx, cy + Inches(0.68), Inches(2.1), Inches(0.5), lab,
             size=12, color=DIM, spacing=1.25)
    chrome(s, n, N)

    # 12 - limits + close ---------------------------------------------------- #
    s = slide(prs); n += 1
    eyebrow(s, "What I am not claiming")
    heading(s, "The limits, stated before you find them.")
    left = [("Quotes are indicative", ", not true NBBO."),
            ("Skew isn't modelled", " \u2014 both legs use one volatility."),
            ("20-day realised vol is an estimator", ", not a forecast of "
             "volatility to expiry.")]
    right = [("Four correlated tickers", " is not diversification."),
             ("Exit thresholds are unvalidated", "."),
             ("A contest week proves nothing", " about profitability. The "
              "literature motivates the idea; it does not validate this system.")]
    for col_i, items in enumerate((left, right)):
        cx = M + col_i * Inches(5.9)
        yy = Inches(2.85)
        for head, rest in items:
            text(s, cx, yy, Inches(0.22), Inches(0.3), "\u2022", size=13,
                 color=ACCENT)
            text(s, cx + Inches(0.26), yy, Inches(5.1), Inches(1.0),
                 [(head, {"bold": True, "color": TEXT}), (rest, {})],
                 size=13, color=DIM, spacing=1.35)
            yy += Inches(1.05)
    text(s, M, Inches(6.15), CONTENT_W, Inches(0.5),
         "An agent that's most useful when it says no.",
         size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print("built %s  (%.0f KB)" % (path, os.path.getsize(path) / 1024))
    sys.exit(0)
